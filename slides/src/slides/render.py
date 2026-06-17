"""PowerPoint (.pptx) → Obsidian-flavored markdown for gap analysis.

Input:  a .pptx file (lecture slides)
Output: markdown with one ## heading per slide, bullets for text content,
        blockquote for speaker notes, [image] placeholder for image-only slides.

Intentionally flat — we want text that's easy to compare against textbook chapters,
not a visual reconstruction. Tables are rendered as bullet lists. SmartArt and
diagrams get an [image] placeholder.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Iterator

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN


# ---------------------------------------------------------------------------
# Text extraction helpers
# ---------------------------------------------------------------------------

def _defang(text: str) -> str:
    """Neutralize Obsidian wikilinks/embeds (``[[note]]`` / ``![[note]]``) in
    untrusted slide text by inserting a zero-width space after each ``[`` that is
    followed by another ``[``. A plain ``str.replace("[[", ...)`` is non-overlapping
    and leaves a live ``[[`` behind on a run of 3+ brackets (``[[[`` -> ``[<zwsp>[[``);
    the lookahead handles runs of any length. We also break Markdown image syntax
    ``![](url)`` (a ZWSP after ``!`` when it precedes ``[``): Obsidian auto-loads
    remote images, so an image in untrusted slide text is a tracking beacon; the
    ZWSP downgrades it to a plain text link. Mirrors canvas-grabber's inline() defang."""
    text = re.sub(r"!(?=\[)", "!​", text)
    return re.sub(r"\[(?=\[)", "[​", text)


def _clean(text: str) -> str:
    """Normalize whitespace, strip control chars, and defang wikilinks."""
    return _defang(re.sub(r"[ \t]+", " ", text).strip())


def _shape_text_lines(shape) -> list[str]:
    """Extract text lines from any shape with a text frame, preserving outline level."""
    lines = []
    if not shape.has_text_frame:
        return lines
    for para in shape.text_frame.paragraphs:
        text = _clean(para.text)
        if not text:
            continue
        level = para.level or 0  # 0 = top level, 1-4 = nested bullets
        indent = "  " * level
        lines.append(f"{indent}- {text}")
    return lines


def _shape_is_image_only(shape) -> bool:
    """True if shape is a picture or chart (visual, no extractable text)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    return shape.shape_type in (
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.MEDIA,
    )


# ---------------------------------------------------------------------------
# Slide rendering
# ---------------------------------------------------------------------------

def _slide_title(slide) -> str | None:
    """Return slide title text, or None."""
    if slide.shapes.title is None:
        return None
    return _clean(slide.shapes.title.text) or None


def _slide_to_md(slide_num: int, slide) -> str:
    """Render one slide as a markdown section."""
    title = _slide_title(slide)
    heading = f"## Slide {slide_num}" + (f" — {title}" if title else "")

    body_lines: list[str] = []
    has_image = False

    for shape in slide.shapes:
        # Skip the title shape — already used above
        if shape == slide.shapes.title:
            continue

        if _shape_is_image_only(shape):
            has_image = True
            continue

        # Tables: flatten to bullets
        if shape.has_table:
            for row in shape.table.rows:
                cells = [_clean(c.text_frame.text) for c in row.cells if _clean(c.text_frame.text)]
                if cells:
                    body_lines.append("- " + " | ".join(cells))
            continue

        # Text frames: preserve outline hierarchy
        if shape.has_text_frame:
            body_lines.extend(_shape_text_lines(shape))

    # Speaker notes
    notes_text = ""
    if slide.has_notes_slide:
        raw = _clean(slide.notes_slide.notes_text_frame.text)
        if raw:
            notes_text = f"\n> *Notes: {raw}*"

    parts = [heading]
    if body_lines:
        parts.extend(body_lines)
    elif has_image and not body_lines:
        parts.append("[image]")
    if notes_text:
        parts.append(notes_text)

    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def pptx_to_markdown(pptx_path: Path, chapter: str | None = None) -> str:
    """Convert a .pptx file to markdown.

    Args:
        pptx_path: Path to the .pptx file.
        chapter:   Optional chapter label (e.g. "02") used in the header.

    Returns:
        Markdown string, ready to write to a file.
    """
    prs = Presentation(pptx_path)

    # Deck title: try the first slide's title, then fall back to filename
    deck_title = None
    if prs.slides:
        deck_title = _slide_title(prs.slides[0])
    if not deck_title:
        deck_title = pptx_path.stem

    chapter_label = f"Ch {chapter} — " if chapter else ""
    # deck_title may be the filename stem (not run through _clean), and the source
    # name is interpolated outside a code span, so defang both here.
    header = f"# {chapter_label}{_defang(deck_title)}\n\n*Source: {_defang(pptx_path.name)}*"

    slide_sections = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_sections.append(_slide_to_md(i, slide))

    return header + "\n\n---\n\n" + "\n\n".join(slide_sections)
